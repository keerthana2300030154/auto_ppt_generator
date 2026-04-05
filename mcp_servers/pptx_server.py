import json
import sys
from typing import Optional
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp import types
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# In-memory store for the active presentation
_prs: Optional[Presentation] = None
_output_path: str = ""

app = Server("pptx-server")


@app.list_tools()
async def list_tools() -> list[types.Tool]:
    return [
        types.Tool(
            name="create_presentation",
            description="Initializes a new blank PowerPoint presentation.",
            inputSchema={
                "type": "object",
                "properties": {
                    "output_path": {
                        "type": "string",
                        "description": "Full file path where the .pptx will be saved"
                    }
                },
                "required": ["output_path"]
            }
        ),
        types.Tool(
            name="add_slide",
            description="Adds a styled slide with a title and bullet points.",
            inputSchema={
                "type": "object",
                "properties": {
                    "title": {"type": "string"},
                    "bullets": {
                        "type": "array",
                        "items": {"type": "string"}
                    },
                    "slide_number": {"type": "integer"},
                    "is_title_slide": {"type": "boolean"}
                },
                "required": ["title", "bullets", "slide_number"]
            }
        ),
        types.Tool(
            name="save_presentation",
            description="Saves the presentation to disk. Call this last.",
            inputSchema={
                "type": "object",
                "properties": {},
                "required": []
            }
        ),
    ]


@app.call_tool()
async def call_tool(name: str, arguments: dict) -> list[types.TextContent]:
    global _prs, _output_path

    if name == "create_presentation":
        _prs = Presentation()
        _prs.slide_width = Inches(13.33)
        _prs.slide_height = Inches(7.5)
        _output_path = arguments["output_path"]
        os.makedirs(os.path.dirname(_output_path), exist_ok=True)
        return [types.TextContent(type="text", text=f"Presentation initialized. Will save to: {_output_path}")]

    elif name == "add_slide":
        if _prs is None:
            return [types.TextContent(type="text", text="ERROR: Call create_presentation first.")]

        title = arguments["title"]
        bullets = arguments.get("bullets", [])
        is_title_slide = arguments.get("is_title_slide", False)

        # Color scheme
        BG_COLOR    = RGBColor(0x1A, 0x1A, 0x2E)
        TITLE_COLOR = RGBColor(0xE9, 0x4F, 0x37)
        BODY_COLOR  = RGBColor(0xF5, 0xF5, 0xF5)
        ACCENT_COLOR = RGBColor(0x16, 0x21, 0x3E)

        W = _prs.slide_width
        H = _prs.slide_height

        blank_layout = _prs.slide_layouts[6]
        slide = _prs.slides.add_slide(blank_layout)

        # Background
        bg = slide.shapes.add_shape(1, 0, 0, W, H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

        if is_title_slide:
            txb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
            tf = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.bold = True
            run.font.size = Pt(48)
            run.font.color.rgb = TITLE_COLOR

            if bullets:
                stxb = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.33), Inches(0.8))
                stf = stxb.text_frame
                sp = stf.paragraphs[0]
                sp.text = bullets[0]
                sp.alignment = PP_ALIGN.CENTER
                sp.runs[0].font.size = Pt(20)
                sp.runs[0].font.color.rgb = BODY_COLOR
        else:
            # Title bar
            title_bar = slide.shapes.add_shape(1, 0, 0, W, Inches(1.2))
            title_bar.fill.solid()
            title_bar.fill.fore_color.rgb = ACCENT_COLOR
            title_bar.line.fill.background()

            txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12.5), Inches(1.0))
            tf = txb.text_frame
            p = tf.paragraphs[0]
            p.text = title
            run = p.runs[0]
            run.font.bold = True
            run.font.size = Pt(32)
            run.font.color.rgb = TITLE_COLOR

            # Bullets
            bullet_top = Inches(1.4)
            bullet_height = Inches(0.65)
            for i, bullet in enumerate(bullets[:5]):
                dot = slide.shapes.add_shape(9, Inches(0.4), bullet_top + i * bullet_height + Inches(0.18), Inches(0.18), Inches(0.18))
                dot.fill.solid()
                dot.fill.fore_color.rgb = TITLE_COLOR
                dot.line.fill.background()

                btxb = slide.shapes.add_textbox(Inches(0.75), bullet_top + i * bullet_height, Inches(12.0), bullet_height)
                btf = btxb.text_frame
                btf.word_wrap = True
                bp = btf.paragraphs[0]
                bp.text = bullet
                bp.runs[0].font.size = Pt(17)
                bp.runs[0].font.color.rgb = BODY_COLOR

        return [types.TextContent(type="text", text=f"Slide added: '{title}'")]

    elif name == "save_presentation":
        if _prs is None:
            return [types.TextContent(type="text", text="ERROR: No presentation to save.")]
        _prs.save(_output_path)
        return [types.TextContent(type="text", text=f"SUCCESS: Presentation saved to {_output_path}")]

    return [types.TextContent(type="text", text=f"Unknown tool: {name}")]


async def main():
    async with stdio_server() as (read_stream, write_stream):
        await app.run(read_stream, write_stream, app.create_initialization_options())

if __name__ == "__main__":
    import asyncio
    asyncio.run(main())