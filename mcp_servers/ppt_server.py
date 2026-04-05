from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt


class PPTServer:
    def __init__(self):
        self.prs = None
        self.background_color = RGBColor(18, 48, 90)
        self.body_color = RGBColor(235, 235, 235)

    def create_presentation(self, filename):
        self.prs = Presentation()
        self.filename = filename

    def _set_slide_background(self, slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.background_color

    def add_title_slide(self, title, subtitle):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        self._set_slide_background(slide)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_slide(self, title, content):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_slide_background(slide)

        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()

        for line in content.split("\n"):
            if line.strip():
                p = tf.add_paragraph()
                p.text = line.strip()
                p.font.size = Pt(18)
                p.font.color.rgb = self.body_color

    def save(self):
        self.prs.save(self.filename)